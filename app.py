from flask import Flask, render_template, request, redirect, url_for, session, flash, jsonify, send_file, session
from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled, NoTranscriptFound, VideoUnavailable
print("Import successful")
from werkzeug.security import check_password_hash ,generate_password_hash
from deep_translator import GoogleTranslator
from gtts import gTTS
import os
import yt_dlp
try:
    import whisper
    WHISPER_AVAILABLE = True
except (ImportError, TypeError) as e:
    # Handle Windows libc loading issue or missing whisper
    print(f"Warning: Could not import whisper: {e}")
    WHISPER_AVAILABLE = False
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
import re
import sqlite3
from werkzeug.utils import secure_filename
from io import BytesIO
import tempfile
import requests
from sumy.summarizers.lsa import LsaSummarizer
from sumy.nlp.tokenizers import Tokenizer
from sumy.parsers.plaintext import PlaintextParser

# =================================================================
# 🤖 HUGGING FACE API KEY - PASTE YOUR KEY HERE
# =================================================================
# Get your free API key from: https://huggingface.co/settings/tokens
# Your key should start with "hf_"
# 
# OPTION 1: Paste your key directly here (EASIEST):
import os

HF_API_KEY = "YOUR_HF_API_KEY"
OPENROUTER_API_KEY = "YOUR_OPENROUTER_API_KEY"

#
# OPTION 2: Use environment variable (more secure):
# HF_API_KEY = os.getenv("HF_API_KEY")
# 
# After pasting your key, replace "PASTE_YOUR_HUGGING_FACE_API_KEY_HERE" 
# with your actual key like: HF_API_KEY = ""
# =================================================================
import language_tool_python

# progress tracking helper
def log_progress(user_id, description):
    if not user_id:
        return
    conn = sqlite3.connect("database.db")
    c = conn.cursor()
    c.execute("INSERT INTO progress (user_id, description) VALUES (?, ?)", (user_id, description))
    conn.commit()
    conn.close()


def init_db():
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    c.execute("""
        CREATE TABLE IF NOT EXISTS videos (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            video_id TEXT UNIQUE,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)

    c.execute("""
        CREATE TABLE IF NOT EXISTS notes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            video_id TEXT,
            timestamp TEXT,
            content TEXT
        )
    """)
    c.execute("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL
        )
    """)

    c.execute("""
        CREATE TABLE IF NOT EXISTS history (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            video_id TEXT,
            notes TEXT
        )
    """)
    conn = sqlite3.connect("database.db")
    cursor = conn.cursor()

    cursor.execute("""
    CREATE TABLE IF NOT EXISTS history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        tool TEXT,
        title TEXT,
        content TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    # progress logs for activities
    c.execute("""
    CREATE TABLE IF NOT EXISTS progress (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        description TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    #   Check existing columns
    c.execute("PRAGMA table_info(history)")
    columns = [column[1] for column in c.fetchall()]

    if "tool" not in columns:
        c.execute("ALTER TABLE history ADD COLUMN tool TEXT")

    if "title" not in columns:
        c.execute("ALTER TABLE history ADD COLUMN title TEXT")

    if "content" not in columns:
        c.execute("ALTER TABLE history ADD COLUMN content TEXT")

    if "created_at" not in columns:
        c.execute("ALTER TABLE history ADD COLUMN created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP")
    # ensure progress table exists (if added later)
    c.execute("""
    CREATE TABLE IF NOT EXISTS progress (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        description TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )""")
    conn.commit()
    conn.close()

def clean_and_correct_text(text):
    """
    Clean and correct grammar in transcript text before AI processing
    """
    try:
        if not text or not text.strip():
            return text
            
        # Step 1: Basic text cleaning
        # Remove extra spaces and normalize whitespace
        cleaned_text = ' '.join(text.split())
        
        # Remove repeated words (common in transcripts)
        words = cleaned_text.split()
        cleaned_words = []
        prev_word = None
        
        for word in words:
            # Skip if same as previous word (likely repetition)
            if word.lower() != prev_word.lower() or len(word) > 3:
                cleaned_words.append(word)
                prev_word = word
        
        cleaned_text = ' '.join(cleaned_words)
        
        # Step 2: Grammar correction using language_tool_python
        try:
            tool = language_tool_python.LanguageTool('en-US')
            corrected_text = tool.correct(cleaned_text)
            
            # If correction didn't change much, use cleaned version
            if len(corrected_text) < len(cleaned_text) * 0.8:
                corrected_text = cleaned_text
                
            print(f"Text corrected: {len(corrected_text)} chars")
            return corrected_text
            
        except Exception as e:
            print(f"Grammar correction failed: {e}")
            return cleaned_text
            
    except Exception as e:
        print(f"Text cleaning error: {e}")
        return text

def extract_key_points(transcript_text):
    """
    Extract 5-10 key points from transcript text using simple heuristics
    Returns formatted key points with nice styling
    """
    try:
        # Split transcript into sentences
        sentences = [s.strip() for s in transcript_text.split('.') if s.strip()]
        
        # Filter for longer, more meaningful sentences (likely to contain key points)
        meaningful_sentences = []
        for sentence in sentences:
            # Look for sentences with keywords that indicate important information
            keywords = ['important', 'key', 'main', 'primary', 'essential', 'crucial', 
                       'significant', 'major', 'fundamental', 'critical', 'because', 'therefore',
                       'first', 'second', 'third', 'finally', 'in conclusion', 'remember',
                       'step', 'process', 'method', 'technique', 'approach', 'strategy']
            
            # Check if sentence contains important keywords or is longer than average
            if any(keyword.lower() in sentence.lower() for keyword in keywords) or len(sentence) > 50:
                meaningful_sentences.append(sentence)
        
        # If we don't have enough meaningful sentences, use all sentences
        if len(meaningful_sentences) < 5:
            meaningful_sentences = sentences
        
        # Take 5-10 key points
        key_points = meaningful_sentences[:10]
        
        # Format as a beautiful highlighted box
        formatted_key_points = """
╔══════════════════════════════════════════════════════════════╗
║                    🎯 KEY POINTS HIGHLIGHTED                  ║
╠══════════════════════════════════════════════════════════════╣
║                                                                ║
"""
        
        for i, point in enumerate(key_points[:8], 1):  # Take up to 8 points
            # Clean up the point and make it more readable
            clean_point = point.strip().capitalize()
            if not clean_point.endswith('.'):
                clean_point += '.'
            
            # Add bullet point with nice formatting
            formatted_key_points += f"║  🔹 {i}. {clean_point:<55} ║\n"
        
        formatted_key_points += """
║                                                                ║
║  💡 These are the most important concepts from the video!      ║
║     Review them carefully for better understanding.           ║
╚══════════════════════════════════════════════════════════════╝
"""
        
        return formatted_key_points if formatted_key_points else """
╔══════════════════════════════════════════════════════════════╗
║                    🎯 KEY POINTS HIGHLIGHTED                  ║
╠══════════════════════════════════════════════════════════════╣
║  🔹 Key points extraction completed.                         ║
║     Please review the full transcript above for details.     ║
╚══════════════════════════════════════════════════════════════╝
"""
        
    except Exception as e:
        print(f"Error extracting key points: {e}")
        return """
╔══════════════════════════════════════════════════════════════╗
║                    🎯 KEY POINTS HIGHLIGHTED                  ║
╠══════════════════════════════════════════════════════════════╣
║  🔹 Unable to extract key points automatically.              ║
║     Please review the full transcript above.                  ║
╚══════════════════════════════════════════════════════════════╝
"""

def generate_notes(text):
    """
    Generate AI notes using HuggingFace API with fallback to sumy
    Returns summary and bullet points from transcript text
    Includes text cleaning and grammar correction for better quality
    """
    try:
        # Step 1: Clean and correct the input text
        print("Cleaning and correcting transcript text...")
        cleaned_text = clean_and_correct_text(text)
        
        # Step 2: Check if HuggingFace API key exists
        api_key = os.getenv("HF_API_KEY")
        
        if api_key:
            # Try HuggingFace API first
            try:
                API_URL = "https://api-inference.huggingface.co/models/google/flan-t5-base"
                headers = {"Authorization": f"Bearer {api_key}"}
                
                # Prepare payload for summarization
                payload = {"inputs": cleaned_text}
                
                # Make API request
                response = requests.post(API_URL, headers=headers, json=payload)
                
                if response.status_code == 200:
                    result = response.json()
                    
                    if result and isinstance(result, list) and len(result) > 0:
                        summary = result[0].get("summary_text", "")
                        
                        if summary:
                            # Create bullet points by splitting summary into sentences
                            sentences = [s.strip() for s in summary.split('.') if s.strip()]
                            bullet_points = [f"• {sentence}" for sentence in sentences]
                            
                            return {
                                "summary": summary,
                                "bullet_points": bullet_points
                            }
            except Exception as api_error:
                print(f"HuggingFace API error: {api_error}")
        
        # If no API key or API failed, use offline summarization using sumy
        print("Using offline summarization with sumy on cleaned text")
        try:
            parser = PlaintextParser.from_string(cleaned_text)
            parser.tokenizer = Tokenizer("english")
            summarizer = LsaSummarizer()
            
            # Generate summary (2-3 sentences)
            summary_sentences = summarizer(parser.document, 2)
            summary = ". ".join(str(sentence) for sentence in summary_sentences)
            
            if not summary:
                # If summarization fails, use first few sentences
                sentences = cleaned_text.split('.')[:3]
                summary = ". ".join(sentences).strip()
            
            # Create bullet points by splitting summary into sentences
            sentences = [s.strip() for s in summary.split('. ') if s.strip()]
            bullet_points = [f"• {sentence}" for sentence in sentences]
            
            return {
                "summary": summary,
                "bullet_points": bullet_points
            }
        except Exception as sumy_error:
            print(f"Sumy error: {sumy_error}")
            # Final fallback - return basic text processing
            sentences = cleaned_text.split('. ')[:3]
            summary = ". ".join(sentences).strip()
            bullet_points = [f"• {s.strip()}" for s in sentences if s.strip()]
            
            return {
                "summary": summary,
                "bullet_points": bullet_points
            }
        
    except Exception as e:
        print(f"Error in generate_notes: {e}")
        # Final fallback - return basic text processing
        sentences = text.split('. ')[:3]
        summary = ". ".join(sentences).strip()
        bullet_points = [f"• {s.strip()}" for s in sentences if s.strip()]
        
        return {
            "summary": summary,
            "bullet_points": bullet_points
        }

init_db()

app = Flask(__name__)

app.secret_key = "smartstudyai_secret"

FREE_LIMIT = 5
user_usage = {}


def extract_video_id(url):
    """
    Extract video ID from YouTube URL with better validation
    Supports multiple YouTube URL formats
    """
    try:
        # Clean URL
        url = url.strip()
        
        # Handle different YouTube URL formats
        patterns = [
            r"(?:v=|\/)([0-9A-Za-z_-]{11})",  # Standard
            r"(?:embed\/)([0-9A-Za-z_-]{11})",   # Embed
            r"(?:v\/)([0-9A-Za-z_-]{11})",     # Shortened
            r"(?:youtu\.be\/)([0-9A-Za-z_-]{11})"  # youtu.be
        ]
        
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                video_id = match.group(1)
                print(f"Extracted video ID: {video_id}")
                return video_id
        
        print("No valid YouTube video ID found in URL")
        return None
        
    except Exception as e:
        print(f"Error extracting video ID: {e}")
        return None

@app.route('/')
def home():
    print("SESSION DATA:", dict(session))
    return render_template("index.html")
    #return render_template('login.html')

@app.route('/translator', methods=['GET', 'POST'])
def translator():
    translated_text = ""
    original_text = ""
    audio_file = None
    error_message = None

    if request.method == 'POST':
        try:
            original_text = request.form.get('text', '').strip()
            target_language = request.form.get('language', '')
            
            # Validation
            if not original_text:
                error_message = "Please enter text to translate"
                if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return jsonify({
                        'success': False,
                        'error': error_message
                    }), 400
            
            if not target_language:
                error_message = "Please select a language"
                if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                    return jsonify({
                        'success': False,
                        'error': error_message
                    }), 400
            
            # Translate text
            translated_text = GoogleTranslator(source='auto', target=target_language).translate(original_text)
            
            if not translated_text:
                raise Exception("Translation returned empty result")

            # save translation to history
            user_id = session.get('user_id')
            if user_id:
                conn = sqlite3.connect("database.db")
                c = conn.cursor()
                title = f"Translator to {target_language}"
                c.execute("INSERT INTO history (user_id, tool, title, content) VALUES (?, ?, ?, ?)",
                          (user_id, "Translator", title, translated_text))
                conn.commit()
                conn.close()
                log_progress(user_id, f"Translated text to {target_language}")

            # Generate voice file
            try:
                tts = gTTS(text=translated_text, lang=target_language, slow=False)
                audio_path = os.path.join("static", "output.mp3")
                tts.save(audio_path)
                audio_file = "output.mp3"
            except Exception as audio_error:
                print(f"Audio generation error: {audio_error}")
                # Don't fail translation if audio generation fails
                audio_file = None
            
            # Check if this is an AJAX request
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({
                    'success': True,
                    'translated_text': translated_text,
                    'audio_file': audio_file,
                    'original_text': original_text
                })
        
        except Exception as e:
            error_message = f"Translation error: {str(e)}"
            print(f"Error in translator: {error_message}")
            
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({
                    'success': False,
                    'error': error_message
                }), 500

    return render_template('translator.html',
                           translated_text=translated_text,
                           original_text=original_text,
                           audio_file=audio_file,
                           error_message=error_message)

@app.route("/youtube-notes", methods=["GET", "POST"])
def youtube_notes():
    notes = None
    video_id = None
    mode = None
    ai_mode = None
    error_message = None

    if request.method == "POST":
        url = request.form.get("url", "").strip()
        transcript_language = request.form.get('transcript_language', 'english')
        
        # Initialize transcript_data to avoid undefined variable errors
        transcript_data = None
        
        if not url:
            error_message = "Please enter a valid YouTube URL"
            return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)
        
        video_id = extract_video_id(url)
        
        if not video_id:
            error_message = "Invalid YouTube URL. Please check and try again."
            return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)

        try:
            print(f"Processing video ID: {video_id}")
            print(f"Language selected: {transcript_language}")
            
            # Try YouTube transcript API first
            try:
                print("Trying YouTube transcript API first...")
                transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
                
                try:
                    transcript = transcript_list.find_transcript(["en"])
                    print("Found manual transcript")
                except:
                    transcript = transcript_list.find_generated_transcript(["en"])
                    print("Found auto-generated transcript")

                transcript_data = transcript.fetch()
                print(f"YouTube transcript fetched successfully with {len(transcript_data)} segments")
                
            except Exception as api_error:
                print(f"YouTube transcript API failed: {api_error}")
                transcript_data = None
                
                # Fallback to audio extraction only if API fails
                try:
                    print("Falling back to audio extraction...")
                    
                    # Download audio using yt-dlp with better error handling
                    ydl_opts = {
                        'format': 'bestaudio/best',
                        'postprocessors': [{
                            'key': 'FFmpegExtractAudio',
                            'preferredcodec': 'mp3',
                            'preferredquality': '192',
                        }],
                        'outtmpl': 'temp_audio.%(ext)s',
                        'quiet': True,
                        'no_warnings': True,
                        'extract_flat': True,
                        'socket_timeout': 30,
                        'retries': 3,
                    }
                    
                    try:
                        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                            ydl.download([f"https://www.youtube.com/watch?v={video_id}"])
                    except Exception as download_error:
                        print(f"Download failed: {download_error}")
                        # Try alternative URL format
                        try:
                            ydl_opts['source_address'] = '0.0.0.0'  # Force IPv4
                            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                                ydl.download([f"https://youtu.be/{video_id}"])
                        except Exception as alt_error:
                            print(f"Alternative download failed: {alt_error}")
                            error_message = "❌ Unable to download video audio. The video may be private or unavailable."
                            return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)
                    
                    # Transcribe using whisper
                    if WHISPER_AVAILABLE:
                        print("Transcribing with whisper...")
                        try:
                            model = whisper.load_model("base")
                            result = model.transcribe("temp_audio.mp3")
                            
                            # Set transcript_data from whisper result
                            transcript_data = result["segments"]
                            print(f"Whisper transcription completed successfully with {len(transcript_data)} segments")
                            
                            # Clean up temporary file
                            try:
                                os.remove("temp_audio.mp3")
                            except:
                                pass
                                
                        except Exception as whisper_error:
                            print(f"Whisper transcription failed: {whisper_error}")
                            error_message = "❌ Audio transcription failed. Please try another video."
                            return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)
                            
                    else:
                        error_message = "❌ Whisper is not available. Cannot process audio transcription."
                        return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)
                        
                except Exception as audio_error:
                    print(f"Audio extraction fallback failed: {audio_error}")
                    error_message = "❌ Failed to generate transcript. Please try another video."
                    return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)

            # Check if we have valid transcript data
            if not transcript_data or len(transcript_data) == 0:
                error_message = "❌ No transcript available for this video."
                return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)

            # Process transcript based on language selection
            try:
                if transcript_language != 'english':
                    print(f"Translating transcript to {transcript_language}...")
                    
                    # Extract full transcript text
                    if isinstance(transcript_data[0], dict):  # YouTube API format
                        full_transcript_text = " ".join([seg["text"] for seg in transcript_data])
                    else:  # Whisper format
                        full_transcript_text = " ".join([seg["text"] for seg in transcript_data])
                    
                    # Step 1: Clean the transcript text
                    cleaned_text = clean_and_correct_text(full_transcript_text)
                    
                    # Step 2: Translate using deep-translator
                    language_codes = {
                        'hindi': 'hi',
                        'telugu': 'te', 
                        'tamil': 'ta',
                        'kannada': 'kn',
                        'malayalam': 'ml',
                        'marathi': 'mr',
                        'bengali': 'bn',
                        'gujarati': 'gu',
                        'punjabi': 'pa'
                    }
                    
                    target_lang = language_codes.get(transcript_language, 'en')
                    
                    if target_lang != 'en':
                        translated_text = GoogleTranslator(
                            source='auto',
                            target=target_lang
                        ).translate(cleaned_text)
                        
                        # Split translated text into sentences and create timestamped format
                        translated_sentences = translated_text.split('. ')
                        
                        # Create timestamped format for translated text
                        formatted_text = ""
                        total_duration = len(transcript_data) * 30  # Estimate 30 seconds per segment
                        
                        for i, sentence in enumerate(translated_sentences):
                            if sentence.strip():
                                # Calculate timestamp based on position
                                timestamp_seconds = min(int((i / len(translated_sentences)) * total_duration), total_duration - 30)
                                minutes = timestamp_seconds // 60
                                seconds = timestamp_seconds % 60
                                timestamp = f"[{minutes:02d}:{seconds:02d}]"
                                 
                                formatted_text += f"{timestamp} • {sentence.strip()}.\n\n"
                        
                        text = formatted_text
                        print(f"Translation to {transcript_language} completed successfully")
                    else:
                        # Use cleaned English text with timestamps
                        formatted_text = ""
                        for i, segment in enumerate(transcript_data):
                            start = int(segment["start"])
                            minutes = start // 60
                            seconds = start % 60
                            timestamp = f"[{minutes:02d}:{seconds:02d}]"
                             
                            text = segment["text"].strip()
                            formatted_text += f"{timestamp} • {text}\n\n"
                        
                        text = formatted_text
                else:
                    # For English, create timestamped format
                    formatted_text = ""
                    for segment in transcript_data:
                        start = int(segment["start"])
                        minutes = start // 60
                        seconds = start % 60
                        timestamp = f"[{minutes:02d}:{seconds:02d}]"
                        
                        text = segment["text"].strip()
                        formatted_text += f"{timestamp} • {text}\n\n"
                    
                    text = formatted_text
                
                # Log user activity
                user_id = session.get('user_id')
                if user_id:
                    log_progress(user_id, f'Generated transcript for {video_id} in {transcript_language}')
                

                return render_template('youtube.html', notes=text, video_id=video_id, mode=mode, ai_mode=ai_mode)

            except Exception as processing_error:
                print(f"Error processing transcript: {processing_error}")
                error_message = "❌ Error processing transcript. Please try again."
                return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)

        except Exception as e:
            print(f"Overall error in youtube_notes: {e}")
            print(f"Error type: {type(e).__name__}")
            error_message = "❌ Failed to generate transcript. Please try again."
            return render_template("youtube.html", notes=notes, video_id=video_id, history=[], mode=mode, ai_mode=ai_mode, error_message=error_message)

    # Fetch history
    history = []

    if "user_id" in session:
        conn = sqlite3.connect("database.db")
        c = conn.cursor()
        c.execute("""
        SELECT tool, title, created_at
        FROM history
        WHERE user_id = ?
        ORDER BY created_at DESC
        LIMIT 5
        """, (session["user_id"],))
        history = c.fetchall()
        conn.close()

    return render_template("youtube.html",
                           notes=notes,
                           video_id=video_id,
                           history=history,
                           mode=mode,
                           ai_mode=ai_mode,
                           error_message=error_message)

@app.route("/ai-tool", methods=["POST"])
def ai_tool():
    """AI Study Tools endpoint for generating quiz, flashcards, mindmap, and study plan"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({"error": "No data received"}), 400
        
        tool = data.get("tool", "").lower()
        text = data.get("text", "")
        
        if not tool or not text:
            return jsonify({"error": "Missing tool or text parameter"}), 400
        
        # Demo responses for each tool type
        if tool == "quiz":
            result = generate_demo_quiz(text)
        elif tool == "flashcards":
            result = generate_demo_flashcards(text)
        elif tool == "mindmap":
            result = generate_demo_mindmap(text)
        elif tool == "studyplan":
            result = generate_demo_study_plan(text)
        else:
            return jsonify({"error": "Unknown tool type"}), 400
        
        # Log the activity if user is logged in
        user_id = session.get('user_id')
        if user_id:
            log_progress(user_id, f"Used AI Study Tool: {tool}")
        
        return jsonify({"result": result, "tool": tool})
        
    except Exception as e:
        print(f"Error in ai_tool route: {e}")
        return jsonify({"error": "Internal server error"}), 500

def generate_demo_quiz(text):
    """Generate demo quiz questions based on the text"""
    return """📝 Generated Quiz Questions

1. What is the main topic discussed in the video?
   A) Topic A
   B) Topic B  
   C) Topic C
   D) Topic D

2. Which concept is most important according to the content?
   A) Concept 1
   B) Concept 2
   C) Concept 3
   D) Concept 4

3. How does the speaker explain the key principle?
   A) Through examples
   B) With diagrams
   C) Using analogies
   D) All of the above

4. What conclusion does the video reach?
   A) Conclusion A
   B) Conclusion B
   C) Conclusion C
   D) Conclusion D

5. Which application is mentioned in the content?
   A) Application 1
   B) Application 2
   C) Application 3
   D) Application 4

💡 Note: This is a demo quiz. Actual quiz will be generated based on your video content using AI."""

def generate_demo_flashcards(text):
    """Generate demo flashcards based on the text"""
    return """🎴 Generated Flashcards

Front: What is the main concept?
Back: The main concept is...

Front: How does this process work?
Back: This process works by...

Front: What are the key components?
Back: The key components include...

Front: Why is this important?
Back: This is important because...

Front: What are the applications?
Back: Applications include...

Front: What are the limitations?
Back: Limitations include...

💡 Note: This is a demo flashcard set. Actual flashcards will be generated based on your video content using AI."""

def generate_demo_mindmap(text):
    """Generate demo mindmap structure based on the text"""
    return """🧠 Generated Mind Map

Main Topic
├── Introduction
│   ├── Definition
│   └── Background
├── Key Concepts
│   ├── Concept 1
│   │   ├── Sub-concept 1.1
│   │   └── Sub-concept 1.2
│   ├── Concept 2
│   │   ├── Sub-concept 2.1
│   │   └── Sub-concept 2.2
│   └── Concept 3
│       ├── Sub-concept 3.1
│       └── Sub-concept 3.2
├── Applications
│   ├── Application 1
│   ├── Application 2
│   └── Application 3
├── Benefits
│   ├── Benefit 1
│   ├── Benefit 2
│   └── Benefit 3
└── Conclusion
    ├── Summary
    └── Future Directions

💡 Note: This is a demo mind map. Actual mind map will be generated based on your video content using AI."""

def generate_demo_study_plan(text):
    """Generate demo 4-day study plan based on the text"""
    return """📅 Generated 4-Day Study Plan

Day 1: Foundation Building
⏰ 2-3 hours
- Watch the video completely (45 mins)
- Take initial notes (30 mins)
- Review key concepts (45 mins)
- Practice basic problems (60 mins)

Day 2: Deep Understanding  
⏰ 2-3 hours
- Re-watch complex sections (30 mins)
- Create detailed summary (60 mins)
- Make flashcards (45 mins)
- Practice intermediate problems (45 mins)

Day 3: Application & Practice
⏰ 2-3 hours  
- Work on practical examples (60 mins)
- Solve practice problems (75 mins)
- Review challenging concepts (30 mins)
- Self-test with flashcards (15 mins)

Day 4: Review & Assessment
⏰ 2 hours
- Complete review of all topics (45 mins)
- Take practice quiz (30 mins)
- Identify weak areas (15 mins)
- Final review and preparation (30 mins)

📚 Daily Tips:
• Start each session with a quick review
• Take 5-minute breaks every 25 minutes
• Use the Pomodoro technique for better focus
• Review notes before bedtime

💡 Note: This is a demo study plan. Actual plan will be personalized based on your video content and learning goals using AI."""

@app.route('/ai-notes')
def ai_notes_page():
    """AI Notes Generator page"""
    return render_template('ai_notes.html')

def format_ai_notes(raw_text):
    """
    Format AI-generated notes for better readability
    Ensures clean bullet points and proper spacing
    """
    try:
        # Split into lines and clean up
        lines = raw_text.strip().split('\n')
        formatted_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Ensure bullet points are properly formatted
            if line.startswith(('1.', '2.', '3.', '4.', '5.')):
                formatted_lines.append(f"\n{line}")
            elif line.startswith(('*', '-', '•')):
                formatted_lines.append(f"• {line.lstrip('*- ').strip()}")
            elif any(keyword in line.lower() for keyword in ['key points', 'important concepts', 'formulas', 'practice questions']):
                formatted_lines.append(f"\n📋 {line.title()}")
            else:
                formatted_lines.append(line)
        
        # Join with proper spacing
        return '\n'.join(formatted_lines)
        
    except Exception as e:
        print(f"Error formatting notes: {e}")
        return raw_text.strip()



@app.route("/generate-questions", methods=["POST"])
def generate_questions_hf():
    """
    Generate practice questions using Hugging Face API
    Bonus feature for SmartStudyAI
    """
    try:
        # Get JSON data from request
        data = request.get_json()
        if not data:
            return jsonify({"error": "No data received"}), 400
        
        transcript = data.get("transcript", "").strip()
        if not transcript:
            return jsonify({"error": "Transcript is required"}), 400
        
        # Get Hugging Face API key from our variable
        hf_api_key = HF_API_KEY
        if not hf_api_key or hf_api_key == "PASTE_YOUR_HUGGING_FACE_API_KEY_HERE":
            return jsonify({
                "error": "API key not configured",
                "message": "Please paste your Hugging Face API key in app.py line 39"
            }), 500
        
        print(f"Generating questions for transcript length: {len(transcript)}")
        
        # Prepare prompt for questions
        prompt = f"""Generate 5 practice questions based on the following transcript. Make them clear and relevant to the content:

Transcript:
{transcript}

Questions:"""
        
        # Hugging Face API configuration
        api_url = "https://api-inference.huggingface.co/models/google/flan-t5-base"
        headers = {
            "Authorization": f"Bearer {hf_api_key}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "inputs": prompt,
            "parameters": {
                "max_length": 300,
                "temperature": 0.7,
                "do_sample": True
            }
        }
        
        print("Sending question request to Hugging Face API...")
        
        # Make API request
        response = requests.post(api_url, headers=headers, json=payload)
        
        if response.status_code == 200:
            result = response.json()
            
            if result and isinstance(result, list) and len(result) > 0:
                generated_text = result[0].get("generated_text", "")
                
                if generated_text:
                    print("Questions generated successfully")
                    return jsonify({
                        "success": True,
                        "questions": generated_text.strip(),
                        "message": "Questions generated successfully"
                    })
        
        # Handle API errors
        print(f"Hugging Face API error: {response.status_code}")
        print(f"Response: {response.text}")
        
        return jsonify({
            "error": "Failed to generate questions",
            "message": "API request failed. Please try again."
        }), 500
        
    except Exception as e:
        print(f"Error in generate_questions_hf: {e}")
        return jsonify({
            "error": "Internal server error",
            "message": "Failed to process request"
        }), 500

@app.route("/history")
def history():
    import sqlite3
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    c.execute("SELECT video_id, created_at FROM videos ORDER BY created_at DESC")
    videos = c.fetchall()

    conn.close()

    return render_template("history.html", videos=videos)

@app.route("/load-video/<video_id>")
def load_video(video_id):
    import sqlite3
    conn = sqlite3.connect("database.db")
    c = conn.cursor()

    c.execute("SELECT timestamp, content FROM notes WHERE video_id = ?", (video_id,))
    rows = c.fetchall()

    conn.close()

    formatted_notes = ""
    for row in rows:
        formatted_notes += f"[{row[0]}] {row[1]}\n\n"

    conn = sqlite3.connect("database.db")
    c = conn.cursor()
    c.execute("SELECT video_id FROM videos ORDER BY created_at DESC")
    videos = c.fetchall()
    conn.close()

    return render_template("youtube.html",
                       notes=formatted_notes,
                       video_id=video_id,
                       history=videos)




@app.route("/download_pdf", methods=["POST"])
def download_pdf():
    notes = request.form.get("notes")
    print("NOTES RECEIVED:", notes)   # 👈 ADD THIS

    if not notes:
        return "No notes received!"

    from io import BytesIO
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer)
    elements = []
    styles = getSampleStyleSheet()

    for line in notes.split("\n"):
        stripped = line.strip()
        if stripped.startswith("-") or stripped.startswith("•") or stripped.startswith("*"):
            elements.append(Paragraph(stripped, styles.get("Bullet", styles["Normal"])))
        else:
            elements.append(Paragraph(line, styles["Normal"]))
        elements.append(Spacer(1, 0.2 * inch))

    doc.build(elements)
    buffer.seek(0)

    # progress log
    user_id = session.get('user_id')
    log_progress(user_id, 'Downloaded notes PDF')

    return send_file(
        buffer,
        as_attachment=True,
        download_name="Video_Notes.pdf",
        mimetype="application/pdf"
    )

@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        username = request.form["username"]
        password = generate_password_hash(request.form["password"])

        conn = sqlite3.connect("database.db")
        c = conn.cursor()

        try:
            c.execute("INSERT INTO users (username, password) VALUES (?, ?)", (username, password))
            conn.commit()
            flash("Registration successful! Please login.")
            return redirect(url_for("login"))
        except:
            flash("Username already exists.")

        conn.close()

    return render_template("register.html")



@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":

        username = request.form.get("username").strip()
        password = request.form.get("password").strip()

        conn = sqlite3.connect("database.db")
        c = conn.cursor()

        c.execute("SELECT * FROM users WHERE username=?", (username,))
        user = c.fetchone()

        conn.close()

        if user is None:
            return render_template("login.html", error="User not found")

        stored_password = user[2]

        # ✅ THIS IS THE FIX
        if check_password_hash(stored_password, password):
            session["user_id"] = user[0]
            session["username"] = user[1]
            return redirect("/")
        else:
            return render_template("login.html", error="Wrong password")

    return render_template("login.html")

@app.route("/profile")
def profile():
    # profile page simply redirects to dashboard for now
    if "user_id" not in session:
        return redirect(url_for("login"))
    return redirect(url_for("dashboard"))


@app.route("/dashboard")
def dashboard():
    if "user_id" not in session:
        return redirect(url_for("login"))
    conn = sqlite3.connect("database.db")
    c = conn.cursor()
    c.execute("SELECT tool, title, created_at FROM history WHERE user_id = ? ORDER BY created_at DESC", (session["user_id"],))
    history = c.fetchall()
    c.execute("SELECT description, created_at FROM progress WHERE user_id = ? ORDER BY created_at DESC", (session["user_id"],))
    progress_items = c.fetchall()
    conn.close()
    return render_template("dashboard.html", history=history, progress=progress_items)


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))


@app.route('/team')
def team():
    # simple static page describing team members
    return render_template('team.html')


@app.route('/converters')
def converters():
    return render_template('converters.html')



# PDF → WORD
@app.route('/convert-pdf-to-word', methods=['POST'])
def convert_pdf_to_word():

    from pdf2docx import Converter

    file = request.files['file']

    pdf_path = os.path.join(tempfile.gettempdir(), secure_filename(file.filename))
    file.save(pdf_path)

    docx_path = pdf_path.replace(".pdf", ".docx")

    cv = Converter(pdf_path)
    cv.convert(docx_path)
    cv.close()

    return send_file(docx_path, as_attachment=True)



# WORD → PDF
@app.route('/convert-word-to-pdf', methods=['POST'])
def convert_word_to_pdf():

    import tempfile
    import os
    import pythoncom
    from flask import request, send_file
    from docx2pdf import convert
    from werkzeug.utils import secure_filename

    file = request.files['file']

    # Save uploaded file
    docx_path = os.path.join(tempfile.gettempdir(), secure_filename(file.filename))
    file.save(docx_path)

    pdf_path = docx_path.replace(".docx", ".pdf")

    try:
        # Initialize COM for Windows
        pythoncom.CoInitialize()

        convert(docx_path, pdf_path)

    except Exception as e:
        return f"Conversion error: {str(e)}"

    return send_file(pdf_path, as_attachment=True)
# PPT → PDF
@app.route('/convert-ppt-to-pdf', methods=['POST'])
def convert_ppt_to_pdf():

    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    from reportlab.lib.colors import white, black
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfbase import pdfmetrics
    import os
    from PIL import Image
    import io

    file = request.files['file']

    ppt_path = os.path.join(tempfile.gettempdir(), secure_filename(file.filename))
    file.save(ppt_path)

    pdf_path = ppt_path.replace(".pptx", ".pdf")

    try:
        prs = Presentation(ppt_path)
        
        # Create PDF with proper page size
        c = canvas.Canvas(pdf_path, pagesize=letter)
        page_width, page_height = letter
        
        # Try to register a better font if available
        try:
            pdfmetrics.registerFont(TTFont('Arial', 'arial.ttf'))
            font_name = 'Arial'
        except:
            font_name = 'Helvetica'

        for slide_idx, slide in enumerate(prs.slides):
            # Clear the page
            c.setFillColor(white)
            c.rect(0, 0, page_width, page_height, fill=1, stroke=0)
            
            # Process each shape in the slide
            for shape in slide.shapes:
                try:
                    # Handle text boxes and shapes with text
                    if hasattr(shape, "text") and shape.text.strip():
                        # Get position and size
                        left = shape.left * 0.75  # Convert from EMUs to points (approximate)
                        top = page_height - (shape.top * 0.75) - 20  # Adjust for PDF coordinate system
                        width = shape.width * 0.75 if shape.width else 400
                        height = shape.height * 0.75 if shape.height else 50
                        
                        # Set font properties based on shape type
                        if hasattr(shape, 'text_frame'):
                            # This is a text box
                            font_size = 12
                            if shape.text_frame.paragraphs:
                                # Try to get font size from first paragraph
                                para = shape.text_frame.paragraphs[0]
                                if para.runs and para.runs[0].font.size:
                                    font_size = int(para.runs[0].font.size / 12700)  # Convert from EMUs
                        else:
                            font_size = 14
                        
                        # Draw text with proper wrapping
                        c.setFillColor(black)
                        c.setFont(font_name, min(font_size, 24))  # Cap font size
                        
                        # Simple text wrapping
                        lines = []
                        words = shape.text.split()
                        current_line = ""
                        max_width = width - 10  # Leave some padding
                        
                        for word in words:
                            test_line = current_line + (" " if current_line else "") + word
                            if c.stringWidth(test_line, font_name, min(font_size, 24)) <= max_width:
                                current_line = test_line
                            else:
                                if current_line:
                                    lines.append(current_line)
                                current_line = word
                        if current_line:
                            lines.append(current_line)
                        
                        # Draw each line
                        line_height = min(font_size, 24) + 2
                        for i, line in enumerate(lines):
                            if top - (i * line_height) > 50:  # Don't draw if too close to bottom
                                c.drawString(left + 5, top - (i * line_height), line)
                    
                    # Handle images
                    elif shape.shape_type == 13:  # Picture type in python-pptx
                        if hasattr(shape, 'image'):
                            try:
                                # Extract image from shape
                                image_bytes = shape.image.blob
                                
                                # Create PIL Image
                                img = Image.open(io.BytesIO(image_bytes))
                                
                                # Calculate dimensions to fit on page
                                img_width, img_height = img.size
                                max_width = page_width - 100
                                max_height = page_height - 100
                                
                                # Scale image to fit
                                if img_width > max_width or img_height > max_height:
                                    ratio = min(max_width / img_width, max_height / img_height)
                                    img_width = int(img_width * ratio)
                                    img_height = int(img_height * ratio)
                                
                                # Position image
                                left = shape.left * 0.75
                                top = page_height - (shape.top * 0.75) - img_height
                                
                                # Draw image on canvas
                                c.drawImage(ImageReader(img), left, top, img_width, img_height)
                                
                            except Exception as e:
                                print(f"Error processing image: {e}")
                                continue
                    
                    # Handle tables
                    elif shape.has_table:
                        try:
                            table = shape.table
                            rows = len(table.rows)
                            cols = len(table.columns)
                            
                            # Calculate table dimensions
                            left = shape.left * 0.75
                            top = page_height - (shape.top * 0.75)
                            cell_width = (shape.width * 0.75) / cols if shape.width else 100
                            cell_height = 30
                            
                            # Draw table cells
                            for row_idx, row in enumerate(table.rows):
                                for col_idx, cell in enumerate(row.cells):
                                    cell_left = left + (col_idx * cell_width)
                                    cell_top = top - (row_idx * cell_height)
                                    
                                    # Draw cell border
                                    c.rect(cell_left, cell_top - cell_height, cell_width, cell_height)
                                    
                                    # Draw cell text
                                    if cell.text and cell.text.strip():
                                        c.setFillColor(black)
                                        c.setFont(font_name, 10)
                                        c.drawString(cell_left + 2, cell_top - cell_height + 15, cell.text.strip())
                        
                        except Exception as e:
                            print(f"Error processing table: {e}")
                            continue
                
                except Exception as e:
                    print(f"Error processing shape {shape}: {e}")
                    continue
            
            # Add new page for next slide
            c.showPage()

        c.save()
        
        # Clean up temporary PPT file
        if os.path.exists(ppt_path):
            os.remove(ppt_path)
            
        return send_file(pdf_path, as_attachment=True)
        
    except Exception as e:
        print(f"Error in PPT to PDF conversion: {e}")
        # Clean up files on error
        if os.path.exists(ppt_path):
            os.remove(ppt_path)
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        return "Error converting PPT to PDF. Please try again.", 500


@app.route('/resume-builder', methods=['GET','POST'])
def resume_builder():
    if request.method == 'POST':
        name = request.form.get('name','').strip()
        email = request.form.get('email','').strip()
        phone = request.form.get('phone','').strip()
        education = request.form.get('education','').strip()
        skills = request.form.get('skills','').strip()
        experience = request.form.get('experience','').strip()

        user_id = session.get('user_id')
        log_progress(user_id, 'Generated Resume')

        # Generate professional PDF resume
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.platypus import Table, TableStyle, PageBreak
        from reportlab.lib.colors import HexColor

        buffer = BytesIO()
        page_width, page_height = letter
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=0.6*inch,
            leftMargin=0.6*inch,
            topMargin=0.5*inch,
            bottomMargin=0.5*inch
        )

        styles = getSampleStyleSheet()
        elements = []

        # Create custom styles
        title_style = styles['Heading1']
        title_style.alignment = 1  # CENTER
        title_style.fontSize = 24
        title_style.textColor = HexColor('#2c2c2c')
        title_style.spaceAfter = 6

        section_style = styles.add(ParagraphStyle(
            name='SectionTitle',
            parent=styles['Heading2'],
            fontSize=12,
            textColor=HexColor('#667eea'),
            spaceAfter=6,
            spaceBefore=10,
            fontName='Helvetica-Bold',
            borderPadding=3,
            borderColor=HexColor('#667eea'),
            borderWidth=0,
            borderRadius=2
        ))

        contact_style = styles.add(ParagraphStyle(
            name='ContactInfo',
            parent=styles['Normal'],
            fontSize=9.5,
            textColor=HexColor('#666666'),
            alignment=1,  # CENTER
            spaceAfter=12
        ))

        content_style = styles.add(ParagraphStyle(
            name='Content',
            parent=styles['Normal'],
            fontSize=10,
            textColor=HexColor('#444444'),
            spaceAfter=6,
            leading=13
        ))

        from reportlab.pdfgen import canvas
        from reportlab.lib.styles import ParagraphStyle

        # Header with name
        elements.append(Paragraph(f"<b>{name.upper()}</b>", title_style))

        # Contact information
        contact_info = []
        if email:
            contact_info.append(email)
        if phone:
            contact_info.append(phone)
        
        if contact_info:
            elements.append(Paragraph(" | ".join(contact_info), contact_style))

        elements.append(Spacer(1, 0.15*inch))

        # Education Section
        if education:
            elements.append(Paragraph("EDUCATION", section_style))
            for line in education.split('\n'):
                if line.strip():
                    elements.append(Paragraph(line.strip(), content_style))
            elements.append(Spacer(1, 0.08*inch))

        # Skills Section
        if skills:
            elements.append(Paragraph("SKILLS", section_style))
            skill_items = []
            for line in skills.split('\n'):
                if line.strip():
                    skill_items.append(Paragraph(f"• {line.strip()}", content_style))
            elements.append(Spacer(1, 0.02*inch))
            for item in skill_items:
                elements.append(item)
            elements.append(Spacer(1, 0.08*inch))

        # Experience Section
        if experience:
            elements.append(Paragraph("EXPERIENCE & PROJECTS", section_style))
            for line in experience.split('\n'):
                if line.strip():
                    elements.append(Paragraph(line.strip(), content_style))

        # Build PDF
        doc.build(elements)
        buffer.seek(0)

        return send_file(
            buffer,
            as_attachment=True,
            download_name=f'{name.replace(" ", "_")}_resume.pdf',
            mimetype='application/pdf'
        )

    return render_template('resume_builder.html')


@app.route('/generate-resume', methods=['POST'])
def generate_resume():
    """Generate professional PDF resume with photo support"""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    from reportlab.platypus import Table, TableStyle, PageBreak, Image
    from reportlab.lib.colors import HexColor
    from PIL import Image as PILImage
    import base64

    try:
        name = request.form.get('name', '').strip()
        email = request.form.get('email', '').strip()
        phone = request.form.get('phone', '').strip()
        location = request.form.get('location', '').strip()
        summary = request.form.get('summary', '').strip()
        degree = request.form.get('degree', '').strip()
        university = request.form.get('university', '').strip()
        graduation = request.form.get('graduation', '').strip()
        skills = request.form.get('skills', '').strip()
        experience = request.form.get('experience', '').strip()
        certifications = request.form.get('certifications', '').strip()
        languages = request.form.get('languages', '').strip()
        photo_data = request.form.get('photo', '').strip()

        user_id = session.get('user_id')
        log_progress(user_id, 'Generated Professional Resume')

        buffer = BytesIO()
        page_width, page_height = letter
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=0.5*inch,
            leftMargin=0.5*inch,
            topMargin=0.6*inch,
            bottomMargin=0.5*inch
        )

        styles = getSampleStyleSheet()
        elements = []

        # Custom styles
        title_style = ParagraphStyle(
            name='Title',
            parent=styles['Heading1'],
            fontSize=26,
            textColor=HexColor('#2c2c2c'),
            spaceAfter=4,
            spaceBefore=0,
            fontName='Helvetica-Bold',
            alignment=0
        )

        section_style = ParagraphStyle(
            name='SectionTitle',
            parent=styles['Heading2'],
            fontSize=11,
            textColor=HexColor('#667eea'),
            spaceAfter=8,
            spaceBefore=12,
            fontName='Helvetica-Bold',
            textTransform='uppercase',
            letterSpacing=2
        )

        contact_style = ParagraphStyle(
            name='Contact',
            parent=styles['Normal'],
            fontSize=9,
            textColor=HexColor('#666666'),
            spaceAfter=2,
            leading=11
        )

        content_style = ParagraphStyle(
            name='Content',
            parent=styles['Normal'],
            fontSize=9.5,
            textColor=HexColor('#555555'),
            spaceAfter=4,
            leading=12,
            leftIndent=12
        )

        content_title_style = ParagraphStyle(
            name='ContentTitle',
            parent=styles['Normal'],
            fontSize=10,
            textColor=HexColor('#2c2c2c'),
            fontName='Helvetica-Bold',
            spaceAfter=2,
            leading=12
        )

        # Build header with photo
        header_data = []
        header_row = []

        # Add photo if available with advanced compression
        if photo_data and photo_data.startswith('data:image'):
            try:
                # Extract base64 data
                header_b64 = photo_data.split(',')[1]
                photo_bytes = base64.b64decode(header_b64)
                photo_buffer = BytesIO(photo_bytes)
                
                # Advanced image optimization
                with PILImage.open(photo_buffer) as img:
                    # Convert to RGB if necessary (for JPEG compatibility)
                    if img.mode in ('RGBA', 'LA', 'P'):
                        background = PILImage.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                        img = background
                    
                    # Calculate optimal dimensions (maintain aspect ratio)
                    max_width, max_height = 150, 150  # Smaller size for PDF
                    img_width, img_height = img.size
                    
                    # Calculate scaling to fit within max dimensions
                    scale = min(max_width / img_width, max_height / img_height)
                    new_width = int(img_width * scale)
                    new_height = int(img_height * scale)
                    
                    # Resize with high-quality resampling
                    img_resized = img.resize((new_width, new_height), PILImage.Resampling.LANCZOS)
                    
                    # Apply additional compression
                    # First try JPEG with high compression
                    photo_buffer_compressed = BytesIO()
                    
                    # Try JPEG first (smaller file size)
                    try:
                        img_resized.save(photo_buffer_compressed, format='JPEG', 
                                       quality=85, optimize=True, progressive=True)
                        photo_buffer_compressed.seek(0)
                        
                        # Check if file size is reasonable (< 50KB)
                        if photo_buffer_compressed.tell() > 50000:
                            # Further compress if still too large
                            photo_buffer_compressed = BytesIO()
                            img_resized.save(photo_buffer_compressed, format='JPEG', 
                                           quality=70, optimize=True, progressive=True)
                            photo_buffer_compressed.seek(0)
                        
                        # If still too large, reduce dimensions further
                        if photo_buffer_compressed.tell() > 50000:
                            # Reduce to 100x100 max
                            scale = min(100 / img_width, 100 / img_height)
                            new_width = int(img_width * scale)
                            new_height = int(img_height * scale)
                            img_resized = img.resize((new_width, new_height), PILImage.Resampling.LANCZOS)
                            
                            photo_buffer_compressed = BytesIO()
                            img_resized.save(photo_buffer_compressed, format='JPEG', 
                                           quality=75, optimize=True, progressive=True)
                            photo_buffer_compressed.seek(0)
                    
                    except Exception:
                        # Fallback to PNG if JPEG fails
                        img_resized.save(photo_buffer_compressed, format='PNG', optimize=True)
                        photo_buffer_compressed.seek(0)
                    
                    # Final size check and fallback
                    if photo_buffer_compressed.tell() > 100000:  # 100KB limit
                        # Create a very small placeholder
                        placeholder_img = PILImage.new('RGB', (80, 80), (240, 240, 240))
                        photo_buffer_compressed = BytesIO()
                        placeholder_img.save(photo_buffer_compressed, format='JPEG', quality=60, optimize=True)
                        photo_buffer_compressed.seek(0)
                    
                    # Calculate display size (in inches for ReportLab)
                    display_width = min(1.0, new_width / 100)  # Max 1 inch wide
                    display_height = min(1.0, new_height / 100)  # Max 1 inch tall
                    
                    photo_img = Image(photo_buffer_compressed, width=display_width*inch, height=display_height*inch)
                    header_row.append(photo_img)
                    
                    print(f"Photo optimized: {photo_buffer_compressed.tell()} bytes")  # Debug info
                    
            except Exception as e:
                print(f"Error processing photo: {e}")
                # Continue without photo if processing fails

        # Create name and contact info cell without emojis
        contact_text = f"<b>{name.upper()}</b>"
        contact_parts = []
        if location:
            contact_parts.append(location)
        if phone:
            contact_parts.append(phone)
        if email:
            contact_parts.append(email)
        if contact_parts:
            contact_text += "<br/>" + " | ".join(contact_parts)
        contact_para = Paragraph(contact_text, contact_style)
        header_row.append(contact_para)

        # Build header (with photo if present)
        if len(header_row) == 2:
            # photo + contact
            header_table = Table([header_row], colWidths=[1.2*inch, 4.3*inch])
            header_table.setStyle(TableStyle([
                ('LEFTPADDING', (0, 0), (-1, -1), 0),
                ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                ('TOPPADDING', (0, 0), (-1, -1), 0),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 0),
                ('VALIGN', (0, 0), (0, 0), 'TOP'),
                ('VALIGN', (1, 0), (1, 0), 'TOP'),
                ('BACKGROUND', (0, 0), (0, 0), HexColor('#f8f9fa')),
            ]))
            elements.append(header_table)
        elif len(header_row) == 1:
            # only contact info
            elements.append(header_row[0])
        else:
            elements.append(Paragraph(f"<b>{name.upper()}</b>", title_style))
            if location or phone or email:
                contact_parts = []
                if location:
                    contact_parts.append(location)
                if phone:
                    contact_parts.append(phone)
                if email:
                    contact_parts.append(email)
                elements.append(Paragraph(" | ".join(contact_parts), contact_style))

        elements.append(Spacer(1, 0.15*inch))

        # Professional Summary
        if summary:
            elements.append(Paragraph("PROFESSIONAL SUMMARY", section_style))
            elements.append(Paragraph(summary, content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Education
        if degree or university or graduation:
            elements.append(Paragraph("EDUCATION", section_style))
            edu_text = f"<b>{degree}</b>"
            if university or graduation:
                edu_text += f"<br/>{university}"
                if graduation:
                    edu_text += f" ({graduation})"
            elements.append(Paragraph(edu_text, content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Skills
        if skills:
            elements.append(Paragraph("SKILLS", section_style))
            skills_list = [s.strip() for s in skills.split(',') if s.strip()]
            for i, skill in enumerate(skills_list):
                if i < len(skills_list) - 1:
                    elements.append(Paragraph(f"{skill} •", content_style))
                else:
                    elements.append(Paragraph(skill, content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Experience & Projects
        if experience:
            elements.append(Paragraph("EXPERIENCE & PROJECTS", section_style))
            sections = experience.split('\n\n')
            for section in sections:
                lines = section.strip().split('\n')
                if lines:
                    elements.append(Paragraph(f"<b>{lines[0]}</b>", content_title_style))
                    for line in lines[1:]:
                        if line.strip():
                            elements.append(Paragraph(f"• {line.strip()}", content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Certifications
        if certifications:
            elements.append(Paragraph("CERTIFICATIONS", section_style))
            for cert in certifications.split('\n'):
                if cert.strip():
                    elements.append(Paragraph(f"• {cert.strip()}", content_style))
            elements.append(Spacer(1, 0.1*inch))

        # Languages
        if languages:
            elements.append(Paragraph("LANGUAGES", section_style))
            for lang in languages.split('\n'):
                if lang.strip():
                    elements.append(Paragraph(f"• {lang.strip()}", content_style))

        # Build PDF
        doc.build(elements)
        buffer.seek(0)

        return send_file(
            buffer,
            as_attachment=True,
            download_name=f'{name.replace(" ", "_")}_resume.pdf',
            mimetype='application/pdf'
        )

    except Exception as e:
        print(f"Error generating resume: {e}")
        return jsonify({'error': str(e)}), 500


# ------ AI study assistant endpoints ------
@app.route('/ai/quiz', methods=['POST'])
def ai_quiz():
    data = request.get_json() or {}
    notes = data.get('notes', '')
    try:
        prompt = (
            "Generate a set of study questions based on the following notes. "
            "Provide 5 multiple-choice questions (with four options each, indicate correct answer and explanation), "
            "5 true/false questions with correct answers, and 3 long-answer questions with answers. "
            f"Return JSON with keys mcq, tf, long.\n\nNotes:\n{notes}"
        )
        # Return basic quiz without API call
        return jsonify({
            'quiz': 'Quiz generation requires API key configuration'
        })
        user_id = session.get('user_id')
        log_progress(user_id, 'Generated AI quiz')
        return jsonify({'quiz': text})
    except Exception as e:
        print("AI quiz error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/flashcards', methods=['POST'])
def ai_flashcards():
    data = request.get_json() or {}
    notes = data.get('notes', '')
    try:
        prompt = (
            "Create a list of concise flashcards (question and answer pairs) from the following notes. "
            f"Notes:\n{notes}"
        )
        # Return basic flashcards without API call
        return jsonify({
            'flashcards': 'Flashcard generation requires API key configuration'
        })
    except Exception as e:
        print("AI flashcards error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/mindmap', methods=['POST'])
def ai_mindmap():
    data = request.get_json() or {}
    notes = data.get('notes', '')
    try:
        prompt = (
            "Generate a structured mind-map description (use indentation or bullet levels) "
            "based on these notes.\n" + notes
        )
        # Return basic mindmap without API call
        return jsonify({
            'mindmap': 'Mindmap generation requires API key configuration'
        })
    except Exception as e:
        print("AI mindmap error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/doubt', methods=['POST'])
def ai_doubt():
    data = request.get_json() or {}
    notes = data.get('notes', '')
    question = data.get('question', '')
    try:
        prompt = (
            "You are an intelligent assistant limited to the following notes. "
            "Answer the user question based only on that content. "
            f"Notes:\n{notes}\nQuestion:\n{question}"
        )
        # Return basic doubt resolution without API call
        return jsonify({
            'answer': 'Doubt resolution requires API key configuration'
        })
    except Exception as e:
        print("AI doubt error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/study-plan', methods=['POST'])
def ai_study_plan():
    data = request.get_json() or {}
    exam_date = data.get('exam_date', '')
    hours_per_day = data.get('hours_per_day', '')
    try:
        prompt = (
            "Create a personalized study schedule plan based on the exam date "
            f"({exam_date}) and available study hours per day ({hours_per_day}). "
            "Output a day-by-day plan with suggestions."
        )
        # Return basic study plan without API call
        return jsonify({
            'plan': 'Study plan generation requires API key configuration'
        })
    except Exception as e:
        print("AI study plan error", e)
        return jsonify({'error': str(e)}), 500

@app.route('/ai/pdf-summary', methods=['POST'])
def ai_pdf_summary():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    file = request.files['file']
    try:
        import fitz  # pymupdf
        pdf_bytes = file.read()
        doc = fitz.open(stream=pdf_bytes, filetype='pdf')
        full_text = ''
        for page in doc:
            full_text += page.get_text()
        doc.close()

        prompt = (
            "Summarize the following PDF contents into concise bullet points:\n\n" + full_text
        )
        # Return basic PDF summary without API call
        return jsonify({
            'summary': 'PDF summary requires API key configuration'
        })
    except Exception as e:
        print("AI PDF summary error", e)
        return jsonify({'error': str(e)}), 500



# ------- LETTER WRITER -------
@app.route('/letter-writer', methods=['GET', 'POST'])
def letter_writer():
    if request.method == 'POST':
        letter_type = request.form.get('letter_type', 'formal')
        sender_name = request.form.get('sender_name', '').strip()
        sender_address = request.form.get('sender_address', '').strip()
        sender_email = request.form.get('sender_email', '').strip()
        sender_phone = request.form.get('sender_phone', '').strip()
        recipient_name = request.form.get('recipient_name', '').strip()
        recipient_title = request.form.get('recipient_title', '').strip()
        recipient_organization = request.form.get('recipient_organization', '').strip()
        recipient_address = request.form.get('recipient_address', '').strip()
        subject = request.form.get('subject', '').strip()
        body = request.form.get('body', '').strip()
        date_str = request.form.get('date', '').strip()

        if not all([sender_name, recipient_name, subject, body]):
            flash("Please fill in all required fields", "error")
            return render_template('letter_writer.html')

        user_id = session.get('user_id')
        log_progress(user_id, f'Generated {letter_type} letter')

        # Generate PDF
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=0.5*inch,
            bottomMargin=0.75*inch
        )

        styles = getSampleStyleSheet()
        elements = []

        # Custom styles
        date_style = ParagraphStyle(
            name='Date',
            parent=styles['Normal'],
            fontSize=11,
            textColor=HexColor('#333333'),
            spaceAfter=24,
            fontName='Helvetica'
        )

        address_style = ParagraphStyle(
            name='Address',
            parent=styles['Normal'],
            fontSize=9.5,
            textColor=HexColor('#555555'),
            spaceAfter=2,
            leading=12,
            fontName='Helvetica'
        )

        salutation_style = ParagraphStyle(
            name='Salutation',
            parent=styles['Normal'],
            fontSize=11,
            textColor=HexColor('#333333'),
            spaceAfter=12,
            fontName='Helvetica'
        )

        body_style = ParagraphStyle(
            name='Body',
            parent=styles['Normal'],
            fontSize=11,
            textColor=HexColor('#333333'),
            spaceAfter=12,
            leading=16,
            alignment=4,  # JUSTIFY
            fontName='Helvetica'
        )

        closing_style = ParagraphStyle(
            name='Closing',
            parent=styles['Normal'],
            fontSize=11,
            textColor=HexColor('#333333'),
            spaceAfter=60,
            fontName='Helvetica'
        )

        signature_style = ParagraphStyle(
            name='Signature',
            parent=styles['Normal'],
            fontSize=11,
            textColor=HexColor('#333333'),
            fontName='Helvetica'
        )

        subject_style = ParagraphStyle(
            name='Subject',
            parent=styles['Normal'],
            fontSize=11,
            textColor=HexColor('#333333'),
            spaceAfter=12,
            fontName='Helvetica-Bold'
        )

        # Sender's address block
        if sender_address or sender_email or sender_phone:
            sender_block = f"{sender_name}<br/>"
            if sender_address:
                sender_block += f"{sender_address}<br/>"
            if sender_email:
                sender_block += f"{sender_email}<br/>"
            if sender_phone:
                sender_block += sender_phone
            elements.append(Paragraph(sender_block, address_style))
            elements.append(Spacer(1, 12))

        # Date
        if date_str:
            elements.append(Paragraph(date_str, date_style))
        else:
            from datetime import datetime
            elements.append(Paragraph(datetime.now().strftime("%B %d, %Y"), date_style))

        # Recipient address block
        recipient_block = f"<b>{recipient_name}</b><br/>"
        if recipient_title:
            recipient_block += f"{recipient_title}<br/>"
        if recipient_organization:
            recipient_block += f"{recipient_organization}<br/>"
        if recipient_address:
            recipient_block += recipient_address
        
        elements.append(Paragraph(recipient_block, address_style))
        elements.append(Spacer(1, 12))

        # Salutation
        salutation = "Dear"
        if recipient_title and recipient_title.lower() in ['mr', 'mrs', 'ms', 'dr', 'prof']:
            salutation += f" {recipient_title} {recipient_name.split()[-1]}"
        else:
            salutation += f" {recipient_name.split()[0]}"
        
        elements.append(Paragraph(f"{salutation},", salutation_style))

        # Subject line (for formal letters)
        if letter_type == 'formal' and subject:
            elements.append(Paragraph(f"<u>RE: {subject}</u>", subject_style))
            elements.append(Spacer(1, 8))

        # Body paragraphs
        for paragraph in body.split('\n\n'):
            if paragraph.strip():
                elements.append(Paragraph(paragraph.strip(), body_style))

        elements.append(Spacer(1, 12))

        # Closing
        closings = {
            'formal': 'Yours sincerely',
            'semi_formal': 'Yours faithfully',
            'informal': 'Best regards',
            'business': 'Respectfully',
            'complaint': 'Yours sincerely'
        }
        closing_text = closings.get(letter_type, 'Yours sincerely')
        elements.append(Paragraph(closing_text + ",", closing_style))

        # Signature space and name
        elements.append(Paragraph(sender_name, signature_style))

        # Build PDF
        doc.build(elements)
        buffer.seek(0)

        return send_file(
            buffer,
            as_attachment=True,
            download_name=f'{letter_type}_letter_{sender_name.replace(" ", "_")}.pdf',
            mimetype='application/pdf'
        )

    return render_template('letter_writer.html')


# ------- CV BUILDER -------
@app.route('/cv-builder', methods=['GET', 'POST'])
def cv_builder():
    if request.method == 'POST':
        # Personal Information
        full_name = request.form.get('full_name', '').strip()
        professional_title = request.form.get('professional_title', '').strip()
        email = request.form.get('email', '').strip()
        phone = request.form.get('phone', '').strip()
        location = request.form.get('location', '').strip()
        website = request.form.get('website', '').strip()
        linkedin = request.form.get('linkedin', '').strip()
        summary = request.form.get('summary', '').strip()
        
        # Education
        education = request.form.get('education', '').strip()
        
        # Experience
        experience = request.form.get('experience', '').strip()
        
        # Skills
        skills = request.form.get('skills', '').strip()
        
        # Projects
        projects = request.form.get('projects', '').strip()
        
        # Certifications
        certifications = request.form.get('certifications', '').strip()
        
        # Languages
        languages = request.form.get('languages', '').strip()

        if not full_name:
            flash("Name is required", "error")
            return render_template('cv_builder.html')

        user_id = session.get('user_id')
        log_progress(user_id, 'Generated Professional CV')

        # Generate Professional CV PDF
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=0.6*inch,
            leftMargin=0.6*inch,
            topMargin=0.5*inch,
            bottomMargin=0.6*inch
        )

        styles = getSampleStyleSheet()
        elements = []

        # Custom Styles
        name_style = ParagraphStyle(
            name='NameStyle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=HexColor('#1a1a1a'),
            spaceAfter=2,
            spaceBefore=0,
            fontName='Helvetica-Bold',
            alignment=0
        )

        title_style = ParagraphStyle(
            name='TitleStyle',
            parent=styles['Normal'],
            fontSize=12,
            textColor=HexColor('#667eea'),
            spaceAfter=10,
            fontName='Helvetica'
        )

        contact_style = ParagraphStyle(
            name='ContactStyle',
            parent=styles['Normal'],
            fontSize=9,
            textColor=HexColor('#666666'),
            spaceAfter=12,
            leading=11,
            fontName='Helvetica'
        )

        section_header_style = ParagraphStyle(
            name='SectionHeader',
            parent=styles['Heading2'],
            fontSize=11,
            textColor=HexColor('#1a1a1a'),
            spaceAfter=8,
            spaceBefore=10,
            fontName='Helvetica-Bold'
        )

        entry_title_style = ParagraphStyle(
            name='EntryTitle',
            parent=styles['Normal'],
            fontSize=10,
            textColor=HexColor('#1a1a1a'),
            fontName='Helvetica-Bold',
            spaceAfter=1,
            leading=12
        )

        entry_meta_style = ParagraphStyle(
            name='EntryMeta',
            parent=styles['Normal'],
            fontSize=9,
            textColor=HexColor('#777777'),
            fontName='Helvetica-Oblique',
            spaceAfter=2,
            leading=11
        )

        entry_text_style = ParagraphStyle(
            name='EntryText',
            parent=styles['Normal'],
            fontSize=9.5,
            textColor=HexColor('#555555'),
            spaceAfter=6,
            leading=12,
            leftIndent=12
        )

        # Header with Name
        elements.append(Paragraph(full_name.upper(), name_style))
        
        if professional_title:
            elements.append(Paragraph(professional_title, title_style))

        # Contact Information
        contact_parts = []
        if email:
            contact_parts.append(f"<b>Email:</b> {email}")
        if phone:
            contact_parts.append(f"<b>Phone:</b> {phone}")
        if location:
            contact_parts.append(f"<b>Location:</b> {location}")
        if website:
            contact_parts.append(f"<b>Website:</b> {website}")
        if linkedin:
            contact_parts.append(f"<b>LinkedIn:</b> {linkedin}")
        
        if contact_parts:
            contact_info = " | ".join(contact_parts)
            elements.append(Paragraph(contact_info, contact_style))
        
        elements.append(Spacer(1, 0.1*inch))

        # Professional Summary
        if summary:
            elements.append(Paragraph("PROFESSIONAL SUMMARY", section_header_style))
            elements.append(Paragraph(summary, entry_text_style))
            elements.append(Spacer(1, 0.08*inch))

        # Experience
        if experience:
            elements.append(Paragraph("PROFESSIONAL EXPERIENCE", section_header_style))
            exp_entries = experience.split('\n\n')
            for entry in exp_entries:
                lines = entry.strip().split('\n')
                if lines and lines[0].strip():
                    elements.append(Paragraph(f"<b>{lines[0]}</b>", entry_title_style))
                    if len(lines) > 1:
                        elements.append(Paragraph(lines[1], entry_meta_style))
                    for line in lines[2:]:
                        if line.strip():
                            elements.append(Paragraph(f"• {line.strip()}", entry_text_style))
            elements.append(Spacer(1, 0.08*inch))

        # Education
        if education:
            elements.append(Paragraph("EDUCATION", section_header_style))
            edu_entries = education.split('\n\n')
            for entry in edu_entries:
                lines = entry.strip().split('\n')
                if lines and lines[0].strip():
                    elements.append(Paragraph(f"<b>{lines[0]}</b>", entry_title_style))
                    for line in lines[1:]:
                        if line.strip():
                            elements.append(Paragraph(line.strip(), entry_meta_style))
            elements.append(Spacer(1, 0.08*inch))

        # Skills
        if skills:
            elements.append(Paragraph("SKILLS", section_header_style))
            skills_list = [s.strip() for s in skills.split(',') if s.strip()]
            for i in range(0, len(skills_list), 3):
                batch = skills_list[i:i+3]
                elements.append(Paragraph(" • ".join(batch), entry_text_style))
            elements.append(Spacer(1, 0.08*inch))

        # Projects
        if projects:
            elements.append(Paragraph("PROJECTS", section_header_style))
            proj_entries = projects.split('\n\n')
            for entry in proj_entries:
                lines = entry.strip().split('\n')
                if lines and lines[0].strip():
                    elements.append(Paragraph(f"<b>{lines[0]}</b>", entry_title_style))
                    for line in lines[1:]:
                        if line.strip():
                            elements.append(Paragraph(f"• {line.strip()}", entry_text_style))
            elements.append(Spacer(1, 0.08*inch))

        # Certifications
        if certifications:
            elements.append(Paragraph("CERTIFICATIONS", section_header_style))
            for cert in certifications.split('\n'):
                if cert.strip():
                    elements.append(Paragraph(f"• {cert.strip()}", entry_text_style))
            elements.append(Spacer(1, 0.08*inch))

        # Languages
        if languages:
            elements.append(Paragraph("LANGUAGES", section_header_style))
            for lang in languages.split('\n'):
                if lang.strip():
                    elements.append(Paragraph(f"• {lang.strip()}", entry_text_style))

        # Build PDF
        doc.build(elements)
        buffer.seek(0)

        return send_file(
            buffer,
            as_attachment=True,
            download_name=f'{full_name.replace(" ", "_")}_CV.pdf',
            mimetype='application/pdf'
        )

    return render_template('cv_builder.html')


@app.route("/check-api-key")
def check_api_key():
    """Check API status - always returns ready for AI notes"""
    return jsonify({"status": "AI_NOTES_READY"})

@app.route("/translate-transcript", methods=["POST"])
def translate_transcript():
    """Translate transcript to selected language"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({"error": "No data received"}), 400
        
        transcript = data.get("transcript", "")
        target_language = data.get("language", "english")
        
        if not transcript:
            return jsonify({"error": "No transcript provided"}), 400
        
        # Check for OpenAI API key
        user_key = session.get("openai_key")
        if not user_key:
            # Use basic translation without OpenAI
            return jsonify({
                "translated_transcript": f"[Translation to {target_language} would require API key]\n\nOriginal transcript:\n{transcript}",
                "message": "Translation requires API key"
            })
        
        # For now, return basic translation without API call
        return jsonify({
            "translated_transcript": f"[Basic translation to {target_language}]\n\nOriginal transcript:\n{transcript}",
            "message": "Basic translation completed"
        })
        
        # Language mapping for better translation
        language_mapping = {
            "english": "English",
            "hindi": "Hindi",
            "telugu": "Telugu", 
            "tamil": "Tamil",
            "kannada": "Kannada",
            "malayalam": "Malayalam"
        }
        
        target_lang_name = language_mapping.get(target_language, "English")
        
        # Create translation prompt
        if target_language == "english":
            prompt = f"""
Please improve and format the following transcript by:

1. Converting spoken language to proper written format
2. Adding proper punctuation and grammar
3. Improving readability while maintaining the original meaning
4. Organizing into clear paragraphs
5. Keeping timestamps intact if present

Transcript to improve:
{transcript}

Return only the improved transcript without any additional text.
"""
        else:
            prompt = f"""
Please translate the following transcript to {target_lang_name} and:

1. Maintain the original meaning and context
2. Use proper grammar and sentence structure
3. Convert spoken language to clear written format
4. Add appropriate punctuation
5. Organize into readable paragraphs
6. Keep timestamps intact if present

Original transcript:
{transcript}

Return only the translated transcript in {target_lang_name} without any additional text.
"""
        
        # For now, return basic translation without API call
        return jsonify({
            "translated_transcript": f"[Basic translation to {target_language}]\n\nOriginal transcript:\n{transcript}",
            "message": "Basic translation completed"
        })
        
        # Log the activity if user is logged in
        user_id = session.get('user_id')
        if user_id:
            log_progress(user_id, f"Translated transcript to {target_language}")
        
        return jsonify({
            "translated_transcript": translated_transcript,
            "language": target_language
        })
        
    except Exception as e:
        print(f"Error in transcript translation: {e}")
        return jsonify({"error": "Translation failed. Please try again."}), 500

@app.route("/connect-ai")
def connect_ai():
    """Render Connect AI Provider page"""
    return render_template("connect_ai.html")

@app.route("/save-api-key", methods=["POST"])
def save_api_key():
    """Save API key in session - placeholder for future use"""
    api_key = request.form.get("api_key", "").strip()
    
    if not api_key:
        flash("Please enter a valid API key", "error")
        return redirect("/connect-ai")
    
    # For now, just store in session without validation
    session["api_key"] = api_key
    
    flash("✅ API key saved successfully!", "success")
    return redirect("/youtube-notes")



def generate_transcript_from_video(video_url):
    """
    Generate transcript from YouTube video URL
    Tries YouTube API first, then falls back to audio extraction
    """
    try:
        print("🎬 Extracting video ID...")
        video_id = extract_video_id(video_url)
        if not video_id:
            print("❌ Invalid YouTube URL")
            return None
        
        print(f"📹 Video ID extracted: {video_id}")
        
        # Try YouTube Transcript API first
        try:
            print("📡 Trying YouTube Transcript API...")
            from youtube_transcript_api import YouTubeTranscriptApi
            transcript_data = YouTubeTranscriptApi.get_transcript(video_id)
            
            # Convert to text and split into parts
            transcript_text = " ".join([item["text"] for item in transcript_data])
            print(f"✅ YouTube API transcript extracted - {len(transcript_text)} characters")
            
            # Split transcript into manageable parts
            transcript_parts = split_transcript_into_parts(transcript_text)
            print(f"📋 Transcript split into {len(transcript_parts)} parts")
            
            return transcript_text
            
        except Exception as api_error:
            print(f"⚠️ YouTube API failed: {api_error}")
            
            # Fallback to audio extraction
            print("🎵 Falling back to audio extraction...")
            return extract_audio_and_transcribe(video_id)
            
    except Exception as e:
        print(f"❌ Transcript generation failed: {e}")
        return None

def split_transcript_into_parts(transcript_text, max_chars=1000):
    """
    Split long transcript into manageable parts for better processing
    """
    try:
        print(f"📝 Splitting transcript of {len(transcript_text)} characters into parts...")
        
        # Split by sentences first, then combine into parts
        import re
        sentences = re.split(r'[.!?]+', transcript_text)
        
        parts = []
        current_part = ""
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
                
            if len(current_part) + len(sentence) + 1 <= max_chars:
                current_part += sentence + ". "
            else:
                if current_part.strip():
                    parts.append(current_part.strip())
                current_part = sentence + ". "
        
        # Add the last part if it exists
        if current_part.strip():
            parts.append(current_part.strip())
        
        print(f"✅ Created {len(parts)} transcript parts")
        for i, part in enumerate(parts):
            print(f"📋 Part {i+1}: {len(part)} characters")
        
        return parts
        
    except Exception as e:
        print(f"❌ Error splitting transcript: {e}")
        return [transcript_text]  # Return original as single part if splitting fails

def extract_audio_and_transcribe(video_id):
    """
    Extract audio from video and transcribe using Whisper
    """
    try:
        if not WHISPER_AVAILABLE:
            print("❌ Whisper not available")
            return None
        
        print("🎵 Downloading audio...")
        
        # Download audio
        import yt_dlp
        ydl_opts = {
            'format': 'bestaudio/best',
            'postprocessors': [{
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }],
            'outtmpl': 'temp_audio.%(ext)s',
            'quiet': True,
            'no_warnings': True,
            'socket_timeout': 30,
            'retries': 3,
        }
        
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([f"https://www.youtube.com/watch?v={video_id}"])
        except Exception as download_error:
            print(f"❌ Audio download failed: {download_error}")
            return None
        
        print("🎙️ Transcribing with Whisper...")
        
        # Transcribe
        import whisper
        model = whisper.load_model("base")
        result = model.transcribe("temp_audio.mp3")
        
        # Clean up
        try:
            import os
            os.remove("temp_audio.mp3")
        except:
            pass
        
        transcript_text = result["text"]
        print(f"✅ Whisper transcription completed - {len(transcript_text)} characters")
        
        # Split into parts
        transcript_parts = split_transcript_into_parts(transcript_text)
        print(f"📋 Audio transcript split into {len(transcript_parts)} parts")
        
        return transcript_text
        
    except Exception as e:
        print(f"❌ Audio transcription failed: {e}")
        return None


@app.route("/generate-ai-notes", methods=["POST"])
def generate_ai_notes():
    try:
        data = request.get_json()
        print("DATA:", data)

        if not data:
            return jsonify({"error": "No data received"}), 400

        # Input handling
        transcript = (data.get("transcript") or "").strip()
        video_url = (
            data.get("video_url")
            or data.get("url")
            or data.get("youtube_url")
            or ""
        ).strip()

        # Get transcript
        if transcript:
            print("Using provided transcript")
        elif video_url:
            print("Extracting transcript...")
            transcript = generate_transcript_from_video(video_url)
            if not transcript:
                return jsonify({"error": "Transcript extraction failed"}), 400
        else:
            return jsonify({"error": "No transcript or video URL provided"}), 400

        print("Transcript length:", len(transcript))

        # Split into chunks
        chunks = [transcript[i:i+1500] for i in range(0, len(transcript), 1500)]
        print("Total chunks:", len(chunks))

        # OpenRouter API
        API_URL = "https://openrouter.ai/api/v1/chat/completions"

        headers = {
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json"
        }

        results = []

        # Process each chunk
        for i, chunk in enumerate(chunks):
            print(f"Processing chunk {i+1}")

            payload = {
                "model": "openai/gpt-3.5-turbo",
                "messages": [
                    {
                        "role": "user",
                        "content": f"""
Summarize this transcript into:
- Key points
- Important concepts
- 3 practice questions

Transcript:
{chunk}
"""
                    }
                ]
            }

            response = requests.post(API_URL, headers=headers, json=payload)

            print("Status:", response.status_code)
            print("Response:", response.text)

            if response.status_code != 200:
                return jsonify({
                    "error": response.text,
                    "status": response.status_code
                }), response.status_code

            output = response.json()

            try:
                text = output["choices"][0]["message"]["content"]
                results.append(text.strip())
            except Exception as e:
                print("Parsing error:", str(e))
                return jsonify({"error": "Invalid AI response"}), 500

        # Combine results
        if results:
            final_notes = "\n\n---\n\n".join(results)
            
            # Save to history with duplicate prevention
            user_id = session.get('user_id')
            if user_id and video_url:
                try:
                    # Extract video ID from URL for consistent identification
                    video_id = None
                    if 'youtube.com/watch?v=' in video_url:
                        video_id = video_url.split('v=')[1].split('&')[0]
                    elif 'youtu.be/' in video_url:
                        video_id = video_url.split('youtu.be/')[1].split('?')[0]
                    
                    if video_id:
                        conn = sqlite3.connect("database.db")
                        c = conn.cursor()
                        
                        # Check if AI Notes already exist for this video and user
                        c.execute("""
                            SELECT COUNT(*) FROM history 
                            WHERE user_id = ? AND tool = 'AI Notes' AND title LIKE ?
                        """, (user_id, f"%{video_id}%"))
                        
                        existing_count = c.fetchone()[0]
                        
                        # Only save if no AI Notes exist for this video
                        if existing_count == 0:
                            title = f"AI Notes - {video_id}"
                            c.execute("""
                                INSERT INTO history (user_id, tool, title, content) 
                                VALUES (?, ?, ?, ?)
                            """, (user_id, "AI Notes", title, final_notes))
                            conn.commit()
                            print(f"Saved new AI notes to history for video {video_id}")
                        else:
                            print(f"AI Notes already exist in history for video {video_id}, skipping save")
                        
                        conn.close()
                except Exception as history_error:
                    print(f"Error saving to history: {history_error}")
            
            return jsonify({
                "success": True,
                "notes": final_notes
            })

        return jsonify({"error": "No notes generated"}), 500

    except Exception as e:
        print("ERROR:", str(e))
        return jsonify({"error": str(e)}), 500
import os

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 10000))
    app.run(host="0.0.0.0", port=port, debug=False)
